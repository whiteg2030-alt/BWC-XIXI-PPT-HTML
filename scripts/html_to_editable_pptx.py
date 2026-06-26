import argparse
import math
import os
import random
import re
import shutil
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5
LIME = RGBColor(183, 255, 0)
WHITE_INK = RGBColor(248, 248, 244)
BLACK_INK = RGBColor(5, 5, 5)
MUTED_BLACK = RGBColor(95, 95, 95)
MUTED_WHITE = RGBColor(150, 150, 145)
GRID_BLACK = RGBColor(24, 24, 24)
GRID_WHITE = RGBColor(224, 224, 218)
BWC_FONT = "白无常可可体"
CHROME_FONT = BWC_FONT
CONTENT_FONT = BWC_FONT
FONT_FILES = {
    "BaiWuchangKeke-Thin.ttf": "白无常可可体 细 (TrueType)",
    "BaiWuchangKeke-Regular.ttf": "白无常可可体 常规 (TrueType)",
    "BaiWuchangKeke-Bold.ttf": "白无常可可体 粗 (TrueType)",
}


def install_packaged_fonts():
    if sys.platform != "win32":
        print("Font install skipped: automatic BaiWuchangKeke installation is Windows-only.")
        return

    font_dir = Path(__file__).resolve().parents[1] / "assets" / "fonts"
    missing = [name for name in FONT_FILES if not (font_dir / name).exists()]
    if missing:
        raise SystemExit(f"Missing packaged font files: {', '.join(missing)}")

    user_fonts = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "Windows" / "Fonts"
    user_fonts.mkdir(parents=True, exist_ok=True)

    installed_any = False
    for file_name, registry_name in FONT_FILES.items():
        src = font_dir / file_name
        dest = user_fonts / file_name
        if not dest.exists() or dest.stat().st_size != src.stat().st_size:
            shutil.copy2(src, dest)
            installed_any = True

        try:
            import winreg

            key_path = r"Software\Microsoft\Windows NT\CurrentVersion\Fonts"
            with winreg.OpenKey(winreg.HKEY_CURRENT_USER, key_path, 0, winreg.KEY_SET_VALUE) as key:
                winreg.SetValueEx(key, registry_name, 0, winreg.REG_SZ, str(dest))
        except Exception as exc:
            print(f"Warning: copied {file_name}, but could not update font registry: {exc}")

    if installed_any:
        print(f"Installed BaiWuchangKeke fonts to {user_fonts}")
    else:
        print("BaiWuchangKeke packaged fonts already available in user font folder.")


def emu(value):
    return Inches(value)


def rgb_from_hex(value):
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def text_of(node, sep="\n"):
    if not node:
        return ""
    text = node.get_text(sep, strip=True)
    return re.sub(r"\n{3,}", "\n\n", text)


def first_text(root, selector, sep="\n"):
    return text_of(root.select_one(selector), sep)


def add_text(slide, text, x, y, w, h, size, color, font=CONTENT_FONT, bold=False, align=None, line_spacing=None):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = str(text).split("\n") if text else [""]
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        if align:
            p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.7)
    else:
        shp.line.fill.background()
    return shp


def add_circle(slide, x, y, d, fill, alpha=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(y), emu(d), emu(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if alpha is not None:
        shp.fill.transparency = alpha
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color, width=0.4):
    shp = slide.shapes.add_connector(1, emu(x1), emu(y1), emu(x2), emu(y2))
    shp.line.color.rgb = color
    shp.line.width = Pt(width)
    return shp


def infer_theme(html, path):
    lower = Path(path).name.lower()
    if "white" in lower or "#f7f7f2" in html:
        return "white"
    return "black"


def palette(theme):
    if theme == "white":
        return {
            "paper": RGBColor(247, 247, 242),
            "ink": BLACK_INK,
            "muted": MUTED_BLACK,
            "grid": GRID_WHITE,
            "panel": RGBColor(245, 245, 240),
            "line": RGBColor(215, 215, 210),
            "particle": RGBColor(72, 72, 72),
        }
    return {
        "paper": RGBColor(0, 0, 0),
        "ink": WHITE_INK,
        "muted": MUTED_WHITE,
        "grid": GRID_BLACK,
        "panel": RGBColor(9, 9, 9),
        "line": RGBColor(36, 36, 36),
        "particle": WHITE_INK,
    }


def draw_background(slide, pal):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["paper"])
    step = 28 / 1672 * SLIDE_W
    x = 0
    while x <= SLIDE_W:
        add_line(slide, x, 0, x, SLIDE_H, pal["grid"], width=0.25)
        x += step
    y = 0
    while y <= SLIDE_H:
        add_line(slide, 0, y, SLIDE_W, y, pal["grid"], width=0.25)
        y += step


def draw_header_footer(slide, pal, page, total, code):
    add_line(slide, 0, 0.66, SLIDE_W, 0.66, pal["line"], width=0.5)
    add_text(slide, "AIXLAB", 0.42, 0.16, 1.2, 0.25, 16, pal["ink"], CHROME_FONT, bold=True)
    add_text(slide, "XIXI STYLE TEMPLATE", 0.43, 0.44, 1.3, 0.15, 4.8, pal["muted"], CHROME_FONT)
    add_text(slide, "01\nTHIN SYSTEM", 2.66, 0.26, 1.2, 0.24, 5.2, pal["ink"], CHROME_FONT, bold=True)
    add_text(slide, "02\nSMALL BEAUTY", 4.68, 0.26, 1.2, 0.24, 5.2, pal["ink"], CHROME_FONT, bold=True)
    add_text(slide, code.upper(), 10.2, 0.29, 0.95, 0.18, 5.2, pal["muted"], CHROME_FONT, align=PP_ALIGN.RIGHT)
    add_text(slide, "2030", 10.66, 0.14, 0.95, 0.35, 22, pal["ink"], CHROME_FONT, bold=True)
    add_rect(slide, 11.88, 0.18, 1.04, 0.31, LIME, radius=True)
    add_text(slide, "→", 12.01, 0.18, 0.3, 0.18, 15, RGBColor(255, 255, 255), CHROME_FONT, bold=True)
    add_text(slide, f"{page:02d}/{total:02d}", 12.36, 0.25, 0.45, 0.12, 5.8, RGBColor(0, 0, 0), CHROME_FONT, bold=True)

    add_text(slide, f"{page:02d}/{total:02d}", 0.42, 7.09, 0.6, 0.13, 5.8, pal["ink"], CHROME_FONT)
    dot_y = 7.13
    start_x = 5.65
    for i in range(total):
        fill = LIME if i < page else pal["paper"]
        dot = add_circle(slide, start_x + i * 0.16, dot_y, 0.06, fill, line=LIME if i < page else pal["muted"])
        dot.line.width = Pt(0.8)
    add_text(slide, f"DESIGN / {code.upper()}", 12.0, 7.08, 0.95, 0.14, 5.2, pal["muted"], CHROME_FONT, align=PP_ALIGN.RIGHT)


def draw_particle_ring(slide, pal, theme):
    random.seed(28 if theme == "black" else 42)
    cx, cy = 9.7, 3.75
    for i in range(260):
        angle = random.random() * math.tau
        radius = random.uniform(1.1, 1.85)
        if 4.95 < angle < 5.55 and random.random() < 0.58:
            continue
        wobble = math.sin(angle * 5.0) * 0.12
        x = cx + math.cos(angle) * (radius + wobble)
        y = cy + math.sin(angle) * (radius * 0.95 + wobble)
        is_lime = random.random() > 0.73
        size = random.uniform(0.018, 0.045) if not is_lime else random.uniform(0.035, 0.08)
        fill = LIME if is_lime else pal["particle"]
        alpha = 0 if is_lime else (22 if theme == "black" else 18)
        add_circle(slide, x, y, size, fill, alpha=alpha)


def draw_cover(slide, section, pal, theme):
    add_text(slide, first_text(section, ".kicker", " "), 0.62, 1.68, 2.7, 0.26, 10, pal["muted"], CHROME_FONT)
    add_line(slide, 0.62, 1.96, 1.18, 1.96, LIME, width=1.6)
    add_text(slide, first_text(section, "h1"), 0.62, 2.28, 3.3, 1.8, 43, pal["ink"], CONTENT_FONT)
    add_text(slide, first_text(section, ".lead", " "), 0.62, 4.73, 5.1, 0.58, 11.5, pal["muted"], CONTENT_FONT, bold=True)
    add_rect(slide, 0.62, 5.68, 1.5, 0.34, LIME, radius=True)
    add_text(slide, first_text(section, ".tag", " "), 0.79, 5.78, 1.15, 0.12, 7.5, RGBColor(0, 0, 0), CHROME_FONT, bold=True)
    draw_particle_ring(slide, pal, theme)


def draw_principle(slide, section, pal):
    add_rect(slide, 0.62, 2.42, 2.88, 2.86, LIME, radius=True)
    num = first_text(section, ".section-num", " ") or "01"
    add_text(slide, num, 1.28, 3.32, 1.5, 0.7, 52, RGBColor(0, 0, 0), CHROME_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, first_text(section, ".kicker", " "), 4.72, 1.42, 2.4, 0.24, 10, pal["muted"], CHROME_FONT)
    add_line(slide, 4.72, 1.69, 5.3, 1.69, LIME, width=1.6)
    add_text(slide, first_text(section, "h1"), 4.72, 2.02, 4.2, 1.6, 40, pal["ink"], CONTENT_FONT)
    add_text(slide, first_text(section, ".lead", " "), 4.74, 4.4, 4.8, 0.34, 11.5, pal["muted"], CONTENT_FONT, bold=True)
    list_items = section.select(".list div")
    y = 5.08
    for item in list_items[:3]:
        b = text_of(item.select_one("b"), " ")
        span = text_of(item.select_one("span"), " ")
        add_circle(slide, 4.75, y + 0.08, 0.07, LIME)
        add_text(slide, b, 4.98, y, 0.3, 0.16, 8.5, LIME, CHROME_FONT, bold=True)
        add_text(slide, span, 5.55, y, 5.2, 0.22, 11.5, pal["muted"], CONTENT_FONT, bold=True)
        add_line(slide, 4.72, y + 0.28, 10.98, y + 0.28, pal["line"], width=0.4)
        y += 0.48


def draw_content(slide, section, pal):
    add_text(slide, first_text(section, ".kicker", " "), 0.62, 1.52, 2.0, 0.2, 8.5, pal["muted"], CHROME_FONT)
    add_line(slide, 0.62, 1.78, 1.08, 1.78, LIME, width=1.4)
    add_text(slide, first_text(section, "h1"), 0.62, 2.04, 4.0, 1.35, 37, pal["ink"], CONTENT_FONT)
    body_nodes = section.select(".body-text")
    y = 4.42
    for body in body_nodes[:2]:
        add_text(slide, text_of(body, " "), 0.62, y, 6.4, 0.42, 10.5, pal["muted"], CONTENT_FONT, bold=True)
        y += 0.52
    add_rect(slide, 9.65, 2.42, 3.05, 1.86, pal["panel"], line=LIME, radius=True)
    metric = section.select_one(".metric-card")
    add_text(slide, first_text(metric, ".num", " ") or "72%", 9.85, 2.72, 1.3, 0.45, 28, LIME, CHROME_FONT, bold=True)
    add_text(slide, first_text(metric, "h3", " "), 9.88, 3.45, 1.5, 0.24, 13, pal["ink"], CONTENT_FONT, bold=True)
    add_line(slide, 9.88, 3.83, 12.45, 3.83, pal["line"], width=0.5)
    add_text(slide, first_text(metric, ".note", " "), 9.88, 4.05, 2.45, 0.3, 8.5, pal["muted"], CONTENT_FONT, bold=True)


def draw_closing(slide, section, pal):
    add_text(slide, first_text(section, ".kicker", " "), 0.62, 1.48, 2.5, 0.22, 8.5, pal["muted"], CHROME_FONT)
    add_line(slide, 0.62, 1.74, 1.12, 1.74, LIME, width=1.4)
    add_text(slide, first_text(section, "h1"), 0.62, 2.03, 3.4, 1.55, 35, pal["ink"], CONTENT_FONT)
    add_text(slide, first_text(section, ".lead", " "), 0.62, 4.56, 4.7, 0.4, 10.5, pal["muted"], CONTENT_FONT, bold=True)
    cards = section.select(".mini-card")
    x = 8.9
    for card in cards[:3]:
        add_rect(slide, x, 2.5, 1.35, 1.2, pal["panel"], line=pal["line"], radius=True)
        add_text(slide, text_of(card.select_one("b"), " "), x + 0.12, 2.62, 0.7, 0.14, 6.5, LIME, CHROME_FONT, bold=True)
        add_text(slide, text_of(card.select_one("h3"), " "), x + 0.12, 2.92, 0.95, 0.25, 9.5, pal["ink"], CONTENT_FONT, bold=True)
        add_text(slide, text_of(card.select_one("p"), " "), x + 0.12, 3.23, 1.05, 0.34, 6.5, pal["muted"], CONTENT_FONT, bold=True)
        x += 1.6


def export_pptx(html_path, output_path):
    html = Path(html_path).read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")
    theme = infer_theme(html, html_path)
    pal = palette(theme)
    sections = soup.select("section.slide-shell")
    if not sections:
        raise SystemExit("No slide-shell sections found.")

    prs = Presentation()
    prs.slide_width = emu(SLIDE_W)
    prs.slide_height = emu(SLIDE_H)
    blank = prs.slide_layouts[6]

    for index, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)
        layout = section.get("data-layout", f"S{index:02d}")
        code = section.get("data-code", layout)
        draw_background(slide, pal)
        draw_header_footer(slide, pal, index, len(sections), code)
        if layout == "S01":
            draw_cover(slide, section, pal, theme)
        elif layout == "S02":
            draw_principle(slide, section, pal)
        elif layout == "S03":
            draw_content(slide, section, pal)
        elif layout == "S04":
            draw_closing(slide, section, pal)
        else:
            add_text(slide, first_text(section, "h1") or f"Slide {index}", 0.7, 1.4, 8.0, 1.0, 34, pal["ink"], CONTENT_FONT)

    prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(description="Export a BWC Xixi HTML deck to editable PPTX.")
    parser.add_argument("html", help="Input HTML deck")
    parser.add_argument("pptx", help="Output editable PPTX")
    args = parser.parse_args()
    install_packaged_fonts()
    export_pptx(args.html, args.pptx)


if __name__ == "__main__":
    main()
